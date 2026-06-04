"""
Build the CTL Agentic AI executive deck (orange theme) from collated /docs content.
Output: docs/CTL_Agentic_AI_Executive_Deck.pptx (16:9).

Audience: C-suite + business + a few architects. Desktop-share demo (not projector),
so use the full canvas, bigger body fonts, minimal blank space.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ── Orange palette ────────────────────────────────────────────────────────────
ORANGE_PRIMARY   = RGBColor(0xE2, 0x6B, 0x0A)
ORANGE_ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE_SOFT      = RGBColor(0xFF, 0xE8, 0xCC)
INK              = RGBColor(0x1F, 0x29, 0x37)
INK_SOFT         = RGBColor(0x4B, 0x55, 0x63)
WHITE            = RGBColor(0xFF, 0xFF, 0xFF)
GREEN            = RGBColor(0x16, 0xA3, 0x4A)
RED_SOFT         = RGBColor(0xDC, 0x26, 0x26)
GREY_BG          = RGBColor(0xF8, 0xFA, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=14, color=INK,
                bullet_color=ORANGE_PRIMARY, line_spacing_pt=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_spacing_pt)
        b = p.add_run(); b.text = "■  "
        b.font.name = "Calibri"; b.font.size = Pt(size)
        b.font.color.rgb = bullet_color; b.font.bold = True
        t = p.add_run(); t.text = item
        t.font.name = "Calibri"; t.font.size = Pt(size); t.font.color.rgb = color

def slide_chrome(slide, title, subtitle=None, page_no=None, total=None):
    add_rect(slide, 0, 0, SW, Inches(0.55), ORANGE_PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.08), Inches(11), Inches(0.4),
             "CTL Agentic AI  ·  Executive Briefing", size=12, bold=True, color=WHITE)
    if page_no:
        add_text(slide, Inches(11.5), Inches(0.08), Inches(1.5), Inches(0.4),
                 f"{page_no} / {total}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.55),
             title, size=28, bold=True, color=INK)
    add_rect(slide, Inches(0.5), Inches(1.27), Inches(0.6), Inches(0.06), ORANGE_ACCENT)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.36), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=INK_SOFT)
    add_rect(slide, 0, SH - Inches(0.18), SW, Inches(0.18), ORANGE_SOFT)

def add_chip(slide, x, y, w, h, text, *, fill=ORANGE_SOFT, color=ORANGE_PRIMARY, size=11):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = color; shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size); run.font.bold = True
    run.font.color.rgb = color

def add_arrow(slide, x, y, w, h, color=ORANGE_PRIMARY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def add_box(slide, x, y, w, h, text, *, fill=WHITE, line=ORANGE_PRIMARY,
            text_color=INK, size=12, bold=True):
    add_rect(slide, x, y, w, h, fill, line=line)
    add_text(slide, x, y, w, h, text, size=size, bold=bold, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

TOTAL = 9

# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, Inches(4.5), SH, ORANGE_PRIMARY)
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), 0, Inches(0.18), SH)
acc.fill.solid(); acc.fill.fore_color.rgb = ORANGE_ACCENT; acc.line.fill.background()

add_text(s, Inches(0.5), Inches(0.7), Inches(3.8), Inches(0.5),
         "CASCADE  ·  CTL", size=15, bold=True, color=ORANGE_SOFT)
add_text(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(2.8),
         "Clear-To-List\nAgentic AI", size=46, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(4.4), Inches(3.8), Inches(2.2),
         ("The judgment step between\nvendor data and the\nlisting decision —\n"
          "automated, grounded,\nand audit-ready."),
         size=18, color=ORANGE_SOFT)
add_text(s, Inches(0.5), Inches(6.7), Inches(3.8), Inches(0.5),
         "Executive Briefing  ·  Business · Architecture · Operations",
         size=12, color=ORANGE_SOFT)

add_text(s, Inches(5.2), Inches(1.2), Inches(7.6), Inches(0.5),
         "The CTL Bottleneck", size=20, bold=True, color=ORANGE_PRIMARY)
add_rect(s, Inches(5.2), Inches(1.8), Inches(0.6), Inches(0.05), ORANGE_ACCENT)
add_text(s, Inches(5.2), Inches(2.0), Inches(7.8), Inches(2.4),
         ("Before any REO, foreclosure, or short-sale\n"
          "asset can be listed, somebody must declare it\n"
          "Clear-To-List — legally clean, properly valued,\n"
          "physically ready — with conditions and citations."),
         size=20, bold=True, color=INK)
add_text(s, Inches(5.2), Inches(4.4), Inches(7.8), Inches(2.0),
         ("Cascade 2.0 already runs the workflow and owns the data.\n"
          "Hardcoded rules can check whether a BPO field is populated\n"
          "and < 90 days old — they can't read title-commitment prose,\n"
          "reconcile valuation disagreements, or apply quarterly\n"
          "investor / state policy. Those calls queue for an analyst,\n"
          "and that's where the days go."),
         size=15, color=INK_SOFT)

add_chip(s, Inches(5.2),  Inches(6.5), Inches(2.4), Inches(0.5), "Policy-cited verdicts", size=12)
add_chip(s, Inches(7.7),  Inches(6.5), Inches(2.4), Inches(0.5), "Independent AI judge", size=12)
add_chip(s, Inches(10.2), Inches(6.5), Inches(2.6), Inches(0.5), "Human-in-the-loop by design", size=12)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "The Problem: a hand-typed verdict is the bottleneck",
             "Cascade 2.0 already routes work and owns data — the decision itself is still manual.",
             page_no=2, total=TOTAL)

col_w = Inches(4.0); gap = Inches(0.15); start_x = Inches(0.5)
top = Inches(2.0); col_h = Inches(4.5)
pains = [
    ("Slow time-to-market",
     ["Every day waiting on CTL is a day of carrying cost — taxes, insurance, HOA, preservation",
      "Routine assets queue behind hard ones — same analyst, same FIFO",
      "Foreclosure-cycle volume spikes blow through capacity"]),
    ("Inconsistent verdicts",
     ["Two analysts on the same asset → two different calls",
      "Same asset reviewed twice in a week can flip",
      "Hard to defend at investor / regulator review"]),
    ("Policy drift goes invisible",
     ["FHA / CWCOT / state / investor rules move quarterly",
      "No system check that the analyst is applying today's rule, not last quarter's",
      "Stale-policy verdicts only surface at audit"]),
]
for i, (h, body) in enumerate(pains):
    x = start_x + i * (col_w + gap)
    add_rect(s, x, top, col_w, Inches(0.55), ORANGE_PRIMARY)
    add_text(s, x, top, col_w, Inches(0.55),
             h, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, top + Inches(0.55), col_w, col_h - Inches(0.55), GREY_BG, line=ORANGE_SOFT)
    add_bullets(s, x + Inches(0.22), top + Inches(0.75), col_w - Inches(0.44), col_h - Inches(0.85),
                body, size=15, color=INK, line_spacing_pt=10)

add_rect(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55), ORANGE_SOFT)
add_text(s, Inches(0.7), Inches(6.82), Inches(12.0), Inches(0.45),
         "What's missing is a JUDGMENT LAYER — facts → defensible verdict with conditions, evidence and policy citations.",
         size=14, bold=True, color=ORANGE_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — Why Agentic vs Rules (C-suite rewrite)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "Why Agentic AI — and not more rules",
             "Rules check whether a field is filled. Agents reason across what's actually in those fields.",
             page_no=3, total=TOTAL)

col_w = Inches(6.0); top = Inches(2.0); col_h = Inches(4.4); gap = Inches(0.3)
# Left — rules
add_rect(s, Inches(0.5), top, col_w, Inches(0.55), RED_SOFT)
add_text(s, Inches(0.5), top, col_w, Inches(0.55),
         "What rules can already do",
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(0.5), top + Inches(0.55), col_w, col_h - Inches(0.55), GREY_BG, line=RED_SOFT)
add_bullets(s, Inches(0.7), top + Inches(0.75), col_w - Inches(0.4), col_h - Inches(0.95),
            ["Check that fields are populated and recent (e.g., BPO < 90 days old)",
             "Flag that two valuations disagree (BPO vs AVM delta > 10%)",
             "Route the case to a fixed analyst queue when a rule fires",
             "Where they stop: cannot read prose, cannot weigh evidence, cannot absorb policy changes without code rewrites"],
            size=15, bullet_color=RED_SOFT, line_spacing_pt=12)

# Right — agents
x2 = Inches(0.5) + col_w + gap
add_rect(s, x2, top, col_w, Inches(0.55), GREEN)
add_text(s, x2, top, col_w, Inches(0.55),
         "What agentic AI adds on top",
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, x2, top + Inches(0.55), col_w, col_h - Inches(0.55), GREY_BG, line=GREEN)
add_bullets(s, x2 + Inches(0.2), top + Inches(0.75), col_w - Inches(0.4), col_h - Inches(0.95),
            ["Reads the free-text vendor narratives — title exceptions, inspector notes, attorney status — the way an analyst would",
             "Decides which valuation to trust and why (comp selection, condition, market signals) — not just that they differ",
             "Applies today's investor / state / program policy via live document lookup — no code change when policy updates",
             "Emits a structured verdict with conditions, evidence trail and citations — analyst-ready, audit-ready"],
            size=15, bullet_color=GREEN, line_spacing_pt=12)

# Bottom band
add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.75), ORANGE_SOFT)
add_text(s, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.35),
         "The split: AI does the judgment. Deterministic code controls the governance.",
         size=14, bold=True, color=ORANGE_PRIMARY)
add_text(s, Inches(0.7), Inches(6.93), Inches(12.0), Inches(0.35),
         "Confidence thresholds, policy enforcement, and human-review escalation live in code — not in the model's hands.",
         size=13, color=INK)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Where the value comes from (ROI / business levers)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "Where the value comes from — six levers",
             "Each lever is a specific place CTL saves money or time. We will put real numbers on each one only after running the agent against real CTL data — not before.",
             page_no=4, total=TOTAL)

levers = [
    ("Faster time-to-list",
     "Routine assets clear in minutes — not days waiting in the analyst queue.",
     "Revenue · Ops"),
    ("Lower carrying cost per asset",
     "Each day off the CTL clock saves taxes, insurance, HOA, preservation and cost-of-capital.",
     "Finance"),
    ("Analyst capacity reclaimed",
     "Senior analysts move off routine throughput and onto ambiguous, high-value cases.",
     "Ops · HR"),
    ("Consistent, defensible verdicts",
     "Same asset → same verdict. Every call is policy-cited and replayable for audit.",
     "Risk · Compliance · Legal"),
    ("Scales with compute, not headcount",
     "Foreclosure-cycle spikes and portfolio acquisitions absorbed without a hiring round.",
     "CFO · COO"),
    ("Policy agility",
     "Quarterly investor / state / program rule changes = update a document, not retrain analysts.",
     "Compliance · Ops"),
]

# 3 columns × 2 rows grid
col_w = Inches(4.05); col_gap = Inches(0.15); start_x = Inches(0.5)
row_h = Inches(2.20); row_gap = Inches(0.18); top = Inches(2.0)

for i, (head, body, owner) in enumerate(levers):
    col = i % 3; row = i // 3
    x = start_x + col * (col_w + col_gap)
    y = top + row * (row_h + row_gap)
    add_rect(s, x, y, col_w, row_h, WHITE, line=ORANGE_PRIMARY)
    # accent strip on the left
    add_rect(s, x, y, Inches(0.10), row_h, ORANGE_PRIMARY)
    # header
    add_text(s, x + Inches(0.25), y + Inches(0.15), col_w - Inches(0.4), Inches(0.55),
             head, size=15, bold=True, color=ORANGE_PRIMARY)
    # body
    add_text(s, x + Inches(0.25), y + Inches(0.75), col_w - Inches(0.4), Inches(1.05),
             body, size=13, color=INK)
    # owner chip
    add_rect(s, x + Inches(0.25), y + row_h - Inches(0.50), col_w - Inches(0.5), Inches(0.35),
             ORANGE_SOFT)
    add_text(s, x + Inches(0.25), y + row_h - Inches(0.50), col_w - Inches(0.5), Inches(0.35),
             owner, size=11, bold=True, color=ORANGE_PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom band — honesty line, taken straight from the source doc
band_y = Inches(6.80)
add_rect(s, Inches(0.5), band_y, Inches(12.3), Inches(0.50), ORANGE_SOFT)
add_text(s, Inches(0.7), band_y, Inches(12.0), Inches(0.50),
         "Today is a working prototype. Each lever gets a real number once we run the agent against a sample of production CTL traffic — those numbers then become the production SLOs and the business case.",
         size=13, bold=True, color=ORANGE_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — How a verdict is produced
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "How a verdict is produced — six steps, every asset",
             "Plan → Investigate (parallel) → Reflect → Apply Policy → Quality Gate → Human Review (only if needed)",
             page_no=5, total=TOTAL)
steps = [
    ("1. Plan",         "Read asset; consult current policy; decide which checks apply for THIS asset"),
    ("2. Investigate",  "Three specialists in parallel — Legal · Valuation · Occupancy — each with scoped data sources"),
    ("3. Reflect",      "Synthesize specialist findings → structured verdict with conditions, evidence and citations"),
    ("4. Apply Policy", "Confidence snapped to discrete buckets; low confidence forced to human review (code, not the LLM)"),
    ("5. Quality Gate", "Independent AI judge scores groundedness vs evidence (1–5). Below threshold → blocked → HITL"),
    ("6. Human Review", "Ambiguous cases reach an analyst with the full evidence package already assembled"),
]
top = Inches(2.0); h = Inches(0.74); w = Inches(12.3); x = Inches(0.5)
for i, (head, body) in enumerate(steps):
    y = top + i * (h + Inches(0.08))
    add_rect(s, x, y, Inches(2.5), h, ORANGE_PRIMARY)
    add_text(s, x + Inches(0.18), y, Inches(2.3), h,
             head, size=17, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x + Inches(2.5), y, w - Inches(2.5), h, GREY_BG, line=ORANGE_SOFT)
    add_text(s, x + Inches(2.7), y, w - Inches(2.85), h,
             body, size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — Solution Architecture (faithful 4-layer diagram)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "Solution Architecture",
             "Four layers — agent workflow, governed tools (MCP), grounded policy (RAG), Azure infrastructure.  Orange-filled blocks use an AI model; outlined blocks are deterministic.",
             page_no=6, total=TOTAL)

fx = Inches(0.4); fy = Inches(1.78); fw = Inches(12.5); fh = Inches(5.50)
# Outer frame — solid, visible boundary around the whole solution
outer = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, fx, fy, fw, fh)
outer.fill.solid(); outer.fill.fore_color.rgb = WHITE
outer.line.color.rgb = ORANGE_PRIMARY; outer.line.width = Pt(2.25)
outer.shadow.inherit = False

# Helper: AI vs non-AI component fill convention
def ai_box(slide, x, y, w, h, text, size=11):
    add_box(slide, x, y, w, h, text, fill=ORANGE_PRIMARY, line=ORANGE_PRIMARY,
            text_color=WHITE, size=size)
def plain_box(slide, x, y, w, h, text, size=11):
    add_box(slide, x, y, w, h, text, fill=WHITE, line=ORANGE_PRIMARY,
            text_color=INK, size=size)

def layer_border(x, y, w, h):
    """Border around a whole layer, no fill."""
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.background()
    shp.line.color.rgb = ORANGE_PRIMARY; shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp

# Geometry — consistent paddings everywhere
LAYER_LEFT     = fx + Inches(0.10)
LAYER_WIDTH    = fw - Inches(0.20)
LAYER_RIGHT    = LAYER_LEFT + LAYER_WIDTH
CONTENT_LEFT   = LAYER_LEFT + Inches(0.20)
CONTENT_WIDTH  = LAYER_WIDTH - Inches(0.40)
HEADER_H       = Inches(0.30)
PAD_HDR_BOX    = Inches(0.12)   # gap between header bar and box row
PAD_BOX_BOTTOM = Inches(0.12)   # gap between last box row and layer border bottom
GAP_BETWEEN_LAYERS = Inches(0.10)
TOP_PAD = Inches(0.12)          # outer frame top → first layer top

def draw_header(y, text):
    add_rect(s, LAYER_LEFT, y, LAYER_WIDTH, HEADER_H, ORANGE_PRIMARY)
    add_text(s, LAYER_LEFT + Inches(0.15), y, LAYER_WIDTH - Inches(0.30), HEADER_H,
             text, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

def centered_start_x(box_w, gap, n):
    total_w = n * box_w + (n - 1) * gap
    return CONTENT_LEFT + (CONTENT_WIDTH - total_w) / 2

# ── Layer 1: Workflow Orchestrator ───────────────────────────────────────────
l1_y = fy + TOP_PAD
draw_header(l1_y, "WORKFLOW ORCHESTRATOR  ·  6-phase decision flow with per-phase timeouts and parallel investigation")

box_h_a = Inches(0.50)
ROW_GAP = Inches(0.10)   # vertical gap between box rows in L1

# Row A: Asset → Planner → Investigation
row_a_y = l1_y + HEADER_H + PAD_HDR_BOX
box_w_a = Inches(2.10); arrow_gap = Inches(0.50)
sx = centered_start_x(box_w_a, arrow_gap, 3)
row_a = [("Asset Profile", False), ("Planner Agent", True), ("Investigation\n(fan-out)", False)]
for i, (t, is_ai) in enumerate(row_a):
    bx = sx + i * (box_w_a + arrow_gap)
    (ai_box if is_ai else plain_box)(s, bx, row_a_y, box_w_a, box_h_a, t, size=12)
    if i < 2:
        add_arrow(s, bx + box_w_a + Inches(0.10), row_a_y + Inches(0.16),
                  Inches(0.30), Inches(0.18))

# Row B: Specialists (all AI)
row_b_y = row_a_y + box_h_a + ROW_GAP
spec_w = Inches(2.20); spec_gap = Inches(0.30)
sx_b = centered_start_x(spec_w, spec_gap, 3)
for i, t in enumerate(["Legal Agent", "Valuation Agent", "Occupancy Agent"]):
    ai_box(s, sx_b + i * (spec_w + spec_gap), row_b_y, spec_w, box_h_a, t, size=12)

# Row C: Reflection (AI) → Policy → Judge (AI) → Human Review → Verdict
row_c_y = row_b_y + box_h_a + ROW_GAP
labs_c = [("Reflection", True),
          ("Policy Enforcer", False),
          ("Quality Gate (Judge)", True),
          ("Human Review", False),
          ("Verdict", False)]
cw = Inches(2.05); cgap = Inches(0.28)
sx_c = centered_start_x(cw, cgap, 5)
for i, (t, is_ai) in enumerate(labs_c):
    x = sx_c + i * (cw + cgap)
    (ai_box if is_ai else plain_box)(s, x, row_c_y, cw, box_h_a, t, size=11)
    if i < len(labs_c) - 1:
        add_arrow(s, x + cw + Inches(0.04), row_c_y + Inches(0.16),
                  Inches(0.20), Inches(0.18))

l1_bottom = row_c_y + box_h_a + PAD_BOX_BOTTOM

# ── Layer 2: MCP Tool Server ─────────────────────────────────────────────────
l2_y = l1_bottom + GAP_BETWEEN_LAYERS
draw_header(l2_y, "MCP TOOL SERVER  ·  single governed entrypoint for every external lookup")
tools = ["Title / HOA / Code Violation", "Occupancy Verification", "AVM / BPO", "Asset Profiler", "Knowledge Base Query"]
tw = Inches(2.20); tgap = Inches(0.18); box_h_t = Inches(0.42)
sx_t = centered_start_x(tw, tgap, 5)
ty = l2_y + HEADER_H + PAD_HDR_BOX
for i, t in enumerate(tools):
    plain_box(s, sx_t + i * (tw + tgap), ty, tw, box_h_t, t, size=11)
l2_bottom = ty + box_h_t + PAD_BOX_BOTTOM

# ── Layer 3: RAG ─────────────────────────────────────────────────────────────
l3_y = l2_bottom + GAP_BETWEEN_LAYERS
draw_header(l3_y, "RAG  ·  policy is content, not code — adding/updating policy never touches orchestration")
rag = [("Policy Knowledge Base", False),
       ("Indexing Pipeline (chunker)", False),
       ("Hybrid Retriever (BM25 + vector)", False),
       ("Embedding Generator", True)]
rw = Inches(2.70); rgap = Inches(0.25); box_h_r = Inches(0.42)
sx_r = centered_start_x(rw, rgap, 4)
ry = l3_y + HEADER_H + PAD_HDR_BOX
for i, (t, is_ai) in enumerate(rag):
    (ai_box if is_ai else plain_box)(s, sx_r + i * (rw + rgap), ry, rw, box_h_r, t, size=11)
l3_bottom = ry + box_h_r + PAD_BOX_BOTTOM

# ── Layer 4: Azure Infrastructure ────────────────────────────────────────────
l4_y = l3_bottom + GAP_BETWEEN_LAYERS
draw_header(l4_y, "AZURE INFRASTRUCTURE  ·  two LLM deployments (Worker + independent Judge), safety, identity, search, observability")
infra = [("AI Foundry\n(Worker LLM)", True),
         ("AI Foundry\n(Judge LLM)", True),
         ("AI Content Safety", True),
         ("AI Search\n(Vector Index)", False),
         ("Entra ID", False),
         ("App Insights", False)]
iw = Inches(1.78); igap = Inches(0.14); box_h_i = Inches(0.46)
sx_i = centered_start_x(iw, igap, 6)
iy = l4_y + HEADER_H + PAD_HDR_BOX
for i, (t, is_ai) in enumerate(infra):
    (ai_box if is_ai else plain_box)(s, sx_i + i * (iw + igap), iy, iw, box_h_i, t, size=10)
l4_bottom = iy + box_h_i + PAD_BOX_BOTTOM

# Layer borders — drawn LAST so they sit over the connectors cleanly
layer_border(LAYER_LEFT, l1_y, LAYER_WIDTH, l1_bottom - l1_y)
layer_border(LAYER_LEFT, l2_y, LAYER_WIDTH, l2_bottom - l2_y)
layer_border(LAYER_LEFT, l3_y, LAYER_WIDTH, l3_bottom - l3_y)
layer_border(LAYER_LEFT, l4_y, LAYER_WIDTH, l4_bottom - l4_y)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — Building Blocks
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "Building blocks — one responsibility each",
             "Replaceable, testable, governable. Adding a vendor or refreshing a policy never touches orchestration.",
             page_no=7, total=TOTAL)
rows = [
    ("Workflow Orchestrator", "Drives 6-phase flow; per-phase timeouts; parallel investigation"),
    ("Planner Agent",         "Reads asset + RAG; decides which checks apply for THIS asset"),
    ("Specialist Agents (3)", "Legal · Valuation · Occupancy — parallel, scoped tools, one domain verdict each"),
    ("Reflection Agent",      "Joins findings; re-grounds against policy; emits asset-level verdict"),
    ("Policy Enforcer",       "Deterministic — schema-validate, snap confidence, force HITL when low (code, not LLM)"),
    ("Quality Gate (Judge)",  "Independent LLM scores groundedness vs evidence; below bar → HITL"),
    ("Human-In-The-Loop",     "Routes flagged cases to analyst with full evidence package; override is captured"),
    ("MCP Tool Server",       "Single governed entrypoint for every external data lookup"),
    ("RAG Layer",             "Policy KB · chunking · hybrid retrieval (BM25 + vector) · embeddings"),
    ("Guardrails Middleware", "Token budget · prompt-injection screening · PII masking on input AND output"),
    ("Resilience Pipelines",  "Polly retries, exponential backoff, per-phase timeouts, circuit-breakers"),
    ("Audit Sink",            "One session id ties every step into one replayable record"),
]
top = Inches(2.0); rh = Inches(0.42); col1 = Inches(3.6); col2 = Inches(8.7)
for i, (a, b) in enumerate(rows):
    y = top + i * rh
    fill = WHITE if i % 2 == 0 else GREY_BG
    add_rect(s, Inches(0.5), y, col1 + col2, rh, fill, line=ORANGE_SOFT)
    add_text(s, Inches(0.65), y, col1 - Inches(0.2), rh,
             a, size=14, bold=True, color=ORANGE_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.65) + col1, y, col2 - Inches(0.2), rh,
             b, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — Tech Stack (AI items first)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "Tech stack — proven, portable, governed",
             "Microsoft-stack first. AI capabilities at the top; platform plumbing below.",
             page_no=8, total=TOTAL)
stack = [
    ("LLM provider",         "Azure AI Foundry — separate Worker + Judge model deployments"),
    ("Agent framework",      "Microsoft Agent Framework — provider-portable agents, tools, and workflows"),
    ("Tools transport",      "Model Context Protocol (MCP) over HTTPS, API-key auth"),
    ("Retrieval (RAG)",      "Azure AI Search — hybrid BM25 + vector (text-embedding-3-small)"),
    ("Safety",               "Azure AI Content Safety + Prompt Shields + Azure Text Analytics PII"),
    ("Evaluation",           "Microsoft.Extensions.AI.Evaluation — Groundedness, Relevance"),
    ("Human-review channel", "Bot Framework — Web Chat & Teams today; channel-portable (Slack, SMS, others)"),
    ("Resilience",           "Polly v8 — retry with exponential backoff, timeout, circuit-breaker"),
    ("Observability",        "OpenTelemetry → Application Insights + per-session JSONL audit"),
    ("Identity",             "Entra ID · DefaultAzureCredential / Managed Identity in production"),
    ("Runtime",              ".NET 9, C#"),
    ("Tests",                "xUnit + NSubstitute · offline regression eval harness"),
]
top = Inches(2.0); rh = Inches(0.42); col1 = Inches(3.6); col2 = Inches(8.7)
for i, (a, b) in enumerate(stack):
    y = top + i * rh
    fill = WHITE if i % 2 == 0 else GREY_BG
    add_rect(s, Inches(0.5), y, col1 + col2, rh, fill, line=ORANGE_SOFT)
    add_text(s, Inches(0.65), y, col1 - Inches(0.2), rh,
             a, size=14, bold=True, color=ORANGE_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.65) + col1, y, col2 - Inches(0.2), rh,
             b, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — Today vs Phase 2 + sequencing
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_chrome(s, "Prototype today  vs  Phase 2  —  what hardens, and in what order",
             "Each gap is a finite engineering item. Vendor + policy + Cascade write-back is the critical path.",
             page_no=9, total=TOTAL)

rows = [
    ("1", "Tool integrations",   "MCP tools with representative mocks",     "Real Title, AVM/BPO, Occupancy, HOA, Code vendors"),
    ("1", "Domain knowledge",    "Sample policy docs",                      "Governed real corpus — named owners, versioning, re-index on change"),
    ("2", "Cascade integration", "Verdict produced as record",              "Camunda trigger → write-back into TaskService"),
    ("3", "Volume",              "Single-asset API runs",                   "Horizontal scale behind a queue; SLOs on production-shape traffic"),
    ("4", "Quality measurement", "Offline regression + runtime judge",      "+ red-team evals, deploy gates, drift monitoring, cost & latency SLOs"),
    ("5", "Region posture",      "Single region",                           "Active-passive multi-region with rehearsed failover"),
    ("6", "LLM provider",        "Azure OpenAI worker + judge",             "Validated against ≥1 alternative — A/B for cost vs quality"),
    ("7", "Compliance",          "Per-decision JSONL + App Insights",       "SOC-2 pack, retention / WORM, residency, model registry"),
    ("8", "Operations",          "Developer-run",                           "Runbooks, on-call, alerting, model upgrade & rollback procedures"),
]
top = Inches(2.0)
c0 = Inches(0.95); c1 = Inches(2.7); c2 = Inches(4.0); c3 = Inches(4.65)
table_w = c0 + c1 + c2 + c3
x0 = Inches(0.5)
hdr_h = Inches(0.5)

add_rect(s, x0, top, c0, hdr_h, INK)
add_text(s, x0, top, c0, hdr_h, "Order",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, x0+c0, top, c1, hdr_h, ORANGE_PRIMARY)
add_text(s, x0+c0+Inches(0.15), top, c1-Inches(0.3), hdr_h, "Area",
         size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, x0+c0+c1, top, c2, hdr_h, ORANGE_PRIMARY)
add_text(s, x0+c0+c1+Inches(0.15), top, c2-Inches(0.3), hdr_h, "Today (prototype)",
         size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, x0+c0+c1+c2, top, c3, hdr_h, ORANGE_PRIMARY)
add_text(s, x0+c0+c1+c2+Inches(0.15), top, c3-Inches(0.3), hdr_h, "Phase 2 (enterprise grade)",
         size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

rh = Inches(0.46)
for i, (o, a, t, p) in enumerate(rows):
    y = top + hdr_h + i*rh
    fill = WHITE if i % 2 == 0 else GREY_BG
    add_rect(s, x0, y, table_w, rh, fill, line=ORANGE_SOFT)
    add_rect(s, x0+Inches(0.18), y+Inches(0.07), Inches(0.6), rh-Inches(0.14), ORANGE_ACCENT)
    add_text(s, x0+Inches(0.18), y+Inches(0.07), Inches(0.6), rh-Inches(0.14),
             o, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0+c0+Inches(0.15), y, c1-Inches(0.3), rh,
             a, size=13, bold=True, color=ORANGE_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0+c0+c1+Inches(0.15), y, c2-Inches(0.3), rh,
             t, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0+c0+c1+c2+Inches(0.15), y, c3-Inches(0.3), rh,
             p, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

band_y = top + hdr_h + len(rows)*rh + Inches(0.15)
add_rect(s, Inches(0.5), band_y, Inches(12.3), Inches(0.55), ORANGE_SOFT)
add_text(s, Inches(0.7), band_y, Inches(12.0), Inches(0.55),
         "Critical path = real vendors + real policy → Cascade write-back → volume benchmark. Everything else hardens that benchmark.",
         size=13, bold=True, color=ORANGE_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)

# ── Save ─────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent.parent / "docs" / "CTL_Agentic_AI_Executive_Deck.pptx"
prs.save(str(out))
print(f"Saved: {out}  ({len(prs.slides)} slides)")
